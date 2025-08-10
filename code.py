import os
import pytesseract
from PIL import Image
from pptx import Presentation
import re
import requests
import json

GEMINI_API_URL = "https://generativelanguage.googleapis.com/v1beta/models/gemini-pro:generateContent"
GEMINI_API_KEY = "Put your api key(i have not put here publically)"

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    all_text = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        all_text.append({"slide_num": idx, "text": "\n".join(slide_text)})
    return all_text

def extract_text_from_images(images_dir):
    all_text = []
    for img_fn in sorted(os.listdir(images_dir)):
        if img_fn.lower().endswith(('.png', '.jpg', '.jpeg')):
            img_path = os.path.join(images_dir, img_fn)
            text = pytesseract.image_to_string(Image.open(img_path))

            match = re.search(r'slide(\d+)', img_fn, re.I)
            slide_num = int(match.group(1)) if match else img_fn
            all_text.append({"slide_num": slide_num, "text": text})
    return all_text

def send_to_gemini(text_blocks):

    slides_desc = "\n\n".join([
        f"Slide {tb['slide_num']}:\n{tb['text']}" for tb in text_blocks
    ])

    prompt = (
        "You are an expert analyst. Given the following multi-slide presentation (slide text below), "
        "identify factual or logical inconsistencies between slides. For every inconsistency (e.g., conflicting revenue figures, contradictory claims, timeline mismatches), "
        "return a JSON array with objects:\n"
        "[\n"
        "{'slides': [slide_numbers], 'issue': description of inconsistency}\n"
        "]\n"
        f"\nPresentation slides:\n{slides_desc}\n"
        "Respond only with the well-structured JSON array. No intro or explanation."
    )

    response = requests.post(
        f"{GEMINI_API_URL}?key={GEMINI_API_KEY}",
        headers={"Content-Type": "application/json"},
        json={"contents": [{"parts": [{"text": prompt}]}]}
    )
    result = response.json()
    try:
        ai_reply = result.get('candidates')[0]['content']['parts'][0]['text']
        issues = json.loads(ai_reply)
        return issues
    except Exception as e:
        print("Gemini output parsing failed:", ai_reply, e)
        return []

def main():
    pptx_file = "slides/sample_deck.pptx"
    images_dir = "slides/images/"

    print("Extracting slide text from PPTX ...")
    ppt_texts = extract_text_from_pptx(pptx_file)
    print("Extracting text from slide images ...")
    img_texts = extract_text_from_images(images_dir)


    slides_by_num = {}
    for tb in ppt_texts + img_texts:
        slides_by_num.setdefault(tb['slide_num'], tb['text'])
    combined_texts = [
        {"slide_num": k, "text": v} for k, v in sorted(slides_by_num.items())
    ]

    print("Sending slide text to Gemini for inconsistency detection ...")
    issues = send_to_gemini(combined_texts)
    print("\n=== Inconsistencies Detected ===")
    for issue in issues:
        slides = ', '.join(str(s) for s in issue.get('slides', []))
        desc = issue.get('issue', '')
        print(f"- Issue on Slide(s): {slides}\n  Description: {desc}\n")
    print("Analysis complete.")

if __name__ == "__main__":
    main()
